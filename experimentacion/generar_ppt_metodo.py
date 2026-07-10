# -*- coding: utf-8 -*-
"""Genera la presentacion del metodo (PPTX) para la fase de experimentacion.
Colores azul + dorado (Universidad San Sebastian). ~10 min, con ejemplo."""

from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = "/home/kpmguser/TesisNicoLLM/Metodo_Gestion_Etica_Priorizacion.pptx"

NAVY  = RGBColor(0x0F, 0x24, 0x47)
NAVY2 = RGBColor(0x1B, 0x3A, 0x66)
GOLD  = RGBColor(0xC9, 0xA2, 0x3E)
GOLDL = RGBColor(0xE9, 0xD7, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK   = RGBColor(0x1B, 0x26, 0x33)
GREY  = RGBColor(0x5B, 0x66, 0x75)
CANVAS= RGBColor(0xF4, 0xF6, 0xF9)
GREEN = RGBColor(0x1E, 0x7A, 0x4F)
RED   = RGBColor(0xB2, 0x3B, 0x3B)
FONT  = "Calibri"

prs = Presentation()
prs.slide_width  = I(13.333)
prs.slide_height = I(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = I(13.333), I(7.5)

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    return s

def rect(s, x, y, w, h, color, line=None, shape=MSO_SHAPE.RECTANGLE):
    r = s.shapes.add_shape(shape, x, y, w, h)
    if color is None: r.fill.background()
    else: r.fill.solid(); r.fill.fore_color.rgb = color
    if line is None: r.line.fill.background()
    else: r.line.color.rgb = line; r.line.width = Pt(1)
    r.shadow.inherit = False
    return r

def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4); tf.margin_top = tf.margin_bottom = Pt(2)
    return tf

def par(tf, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT, before=0, after=4,
        bullet=False, italic=False, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_before = Pt(before); p.space_after = Pt(after)
    if bullet: text = "•   " + text
    r = p.add_run(); r.text = text
    f = r.font; f.name = FONT; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    return p

def runs(tf, parts, size, align=PP_ALIGN.LEFT, before=0, after=4, bullet=False, first=False):
    """parts: lista de (texto, color, bold)."""
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_before = Pt(before); p.space_after = Pt(after)
    if bullet:
        r = p.add_run(); r.text = "•   "; r.font.name = FONT; r.font.size = Pt(size)
        r.font.color.rgb = GOLD; r.font.bold = True
    for txt, color, bold in parts:
        r = p.add_run(); r.text = txt
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return p

def header(s, kicker, title):
    rect(s, 0, 0, SW, I(1.25), NAVY)
    rect(s, 0, I(1.25), SW, Pt(4), GOLD)
    tf = textbox(s, I(0.55), I(0.16), I(12.2), I(1.0), MSO_ANCHOR.MIDDLE)
    par(tf, kicker, 11, GOLD, True, after=1, first=True)
    par(tf, title, 25, WHITE, True, after=0)

def footer(s, n):
    tf = textbox(s, I(0.55), I(7.02), I(9), I(0.35))
    par(tf, "Método de gestión ética y priorización de requisitos con IA", 8.5, GREY, first=True)
    tf2 = textbox(s, I(11.8), I(7.02), I(1.1), I(0.35))
    par(tf2, str(n), 9, GOLD, True, align=PP_ALIGN.RIGHT, first=True)

def style_table(tbl):
    # quita el estilo por defecto con bandas
    tblPr = tbl._tbl.tblPr
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')

def set_cell(cell, text, size=11, color=INK, bold=False, fill=WHITE, align=PP_ALIGN.LEFT):
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Pt(7)
    cell.margin_top = cell.margin_bottom = Pt(4)
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

# ==========================================================================
# 1 — Portada
# ==========================================================================
s = slide(NAVY)
rect(s, 0, I(2.55), SW, Pt(3), GOLD)
rect(s, I(0.9), I(2.62), I(1.4), Pt(6), GOLD)  # detalle
tf = textbox(s, I(0.9), I(2.75), I(11.5), I(2.4))
par(tf, "FASE DE EXPERIMENTACIÓN", 14, GOLD, True, after=6, first=True)
par(tf, "Gestión ética y priorización\nde requisitos con IA", 40, WHITE, True, after=8)
par(tf, "Un método de 8 fases para decidir qué construir — y hacerlo bien.", 16, GOLDL, False)
tf2 = textbox(s, I(0.9), I(6.35), I(11.5), I(0.6))
par(tf2, "Tesis · Nicolás Pérez   |   Universidad San Sebastián   |   Sesión de ~10 minutos", 12, RGBColor(0xB9,0xC6,0xD8), first=True)

# ==========================================================================
# 2 — Objetivo
# ==========================================================================
s = slide(); header(s, "POR QUÉ ESTAMOS AQUÍ", "Qué van a probar hoy")
footer(s, 2)
tf = textbox(s, I(0.7), I(1.75), I(11.9), I(1.4))
par(tf, "Una herramienta que ayuda a los equipos de software a revisar sus requisitos "
        "antes de construirlos: detecta riesgos éticos, propone cómo tratarlos y ordena "
        "los requisitos por su valor real.", 17, INK, first=True, after=6)
# tres tarjetas
cards = [
    ("Detectar", "Identifica tensiones éticas que a simple vista no se ven (sesgo, privacidad, opacidad)."),
    ("Tratar", "Propone reformular, mitigar, aceptar o eliminar, con apoyo en normativa real."),
    ("Priorizar", "Calcula un puntaje que equilibra beneficio y riesgo, y ordena los requisitos."),
]
x = I(0.7); cw = I(3.95); gap = I(0.18)
for i,(t,b) in enumerate(cards):
    cx = x + i*(cw+gap)
    rect(s, cx, I(3.5), cw, I(2.7), CANVAS, line=GOLDL)
    rect(s, cx, I(3.5), cw, Pt(6), GOLD)
    ctf = textbox(s, cx+I(0.25), I(3.8), cw-I(0.5), I(2.2))
    par(ctf, str(i+1), 13, GOLD, True, first=True, after=2)
    par(ctf, t, 20, NAVY, True, after=6)
    par(ctf, b, 13, GREY)

# ==========================================================================
# 3 — El método en 8 fases
# ==========================================================================
s = slide(); header(s, "VISIÓN GENERAL", "El método, en cuatro momentos")
footer(s, 3)
groups = [
    ("Fase 1", "Registro", "Se cargan el proyecto y sus requisitos."),
    ("Fases 2–3", "Análisis ético con IA", "La IA analiza cada requisito y propone tratamientos."),
    ("Fases 4–6", "Priorización", "Se ponderan beneficios y riesgos → ranking."),
    ("Fases 7–8", "Trazabilidad y tablero", "Derivados, bandera ética y exportación."),
]
x = I(0.6); cw = I(2.85); gap = I(0.33); y = I(2.6); ch = I(2.9)
for i,(k,t,b) in enumerate(groups):
    cx = x + i*(cw+gap)
    rect(s, cx, y, cw, ch, NAVY if i%2==0 else NAVY2)
    rect(s, cx, y, cw, Pt(6), GOLD)
    ctf = textbox(s, cx+I(0.22), y+I(0.28), cw-I(0.44), ch-I(0.5))
    par(ctf, k, 12, GOLD, True, first=True, after=3)
    par(ctf, t, 17, WHITE, True, after=8)
    par(ctf, b, 12, RGBColor(0xC7,0xD3,0xE3))
    if i < 3:
        ar = rect(s, cx+cw+Pt(2), y+ch/2-I(0.18), gap-Pt(4), I(0.36), GOLD, shape=MSO_SHAPE.CHEVRON)
tf = textbox(s, I(0.6), I(5.85), I(12.1), I(0.9))
runs(tf, [("Hoy recorreremos las cuatro etapas de principio a fin. ", INK, False),
          ("El corazón del método son las Fases 2–6.", NAVY, True)], 15, first=True)

# ==========================================================================
# 4 — Fases 2-3 análisis ético
# ==========================================================================
s = slide(); header(s, "FASES 2–3", "Análisis ético asistido por IA")
footer(s, 4)
tf = textbox(s, I(0.7), I(1.7), I(11.9), I(0.7))
par(tf, "Por cada requisito, la IA produce un análisis en tres capas — y el humano lo revisa y edita.",
    16, INK, first=True)
layers = [
    ("Capa 1 · Identificación", "Tema ético, actor afectado, tipo de daño y norma tensionada, con "
     "citas reales a la normativa (EU AI Act, GDPR, leyes chilenas 19.628 y 20.609)."),
    ("Capa 2 · Análisis", "Mapa de stakeholders y tensiones entre valores (p. ej. utilidad vs. equidad)."),
    ("Capa 3 · Deliberación", "Opciones de tratamiento, reformulaciones, requisitos derivados y "
     "preguntas para decidir."),
]
y = I(2.6)
for i,(t,b) in enumerate(layers):
    rect(s, I(0.7), y, I(11.9), I(1.15), CANVAS, line=GOLDL)
    rect(s, I(0.7), y, I(0.14), I(1.15), GOLD)
    ntf = textbox(s, I(1.05), y+I(0.13), I(11.3), I(0.95), MSO_ANCHOR.MIDDLE)
    runs(ntf, [(t+"   ", NAVY, True), (b, GREY, False)], 14, first=True)
    y = y + I(1.32)
tf = textbox(s, I(0.7), I(6.55), I(11.9), I(0.5))
runs(tf, [("Decisiones posibles: ", INK, True),
          ("aceptar · reformular · mitigar · eliminar.", NAVY, False)], 14, first=True)

# ==========================================================================
# 5 — Fases 4-6 priorización (fórmula + escalas)
# ==========================================================================
s = slide(); header(s, "FASES 4–6", "Priorización: cómo se calcula el puntaje")
footer(s, 5)
tf = textbox(s, I(0.7), I(1.65), I(6.0), I(4.4))
par(tf, "Cada requisito se puntúa según cuánto aporta y cuánto arriesga.", 15, INK, first=True, after=8)
runs(tf, [("Dimensiones", NAVY, True), (" — criterios con un ", INK, False), ("tipo", INK, True),
          (" y un ", INK, False), ("peso 1–5.", INK, True)], 14, bullet=True, after=6)
runs(tf, [("Fuerza 0–5", NAVY, True), (" — cuánto activa el requisito esa dimensión (0 = no aplica).", INK, False)],
     14, bullet=True, after=10)
runs(tf, [("Suman:  ", INK, True), ("beneficio", GREEN, True), (" y ", INK, False), ("valor ético.", GREEN, True)],
     14, bullet=True, after=4)
runs(tf, [("Restan:  ", INK, True), ("costo", RED, True), (" y ", INK, False), ("riesgo ético.", RED, True)],
     14, bullet=True, after=8)
runs(tf, [("La ética mueve el ranking: ", NAVY, True),
          ("un requisito muy útil pero riesgoso baja frente a otro más seguro.", INK, False)], 14, after=0)
# caja fórmula a la derecha
rect(s, I(7.0), I(2.0), I(5.6), I(3.9), NAVY)
rect(s, I(7.0), I(2.0), I(5.6), Pt(6), GOLD)
ftf = textbox(s, I(7.35), I(2.35), I(4.95), I(3.3), MSO_ANCHOR.MIDDLE)
par(ftf, "PUNTAJE FINAL", 13, GOLD, True, first=True, after=10)
runs(ftf, [("Σ (peso × fuerza) beneficio", WHITE, True)], 15, after=4)
runs(ftf, [("+  Σ (peso × fuerza) valor ético", GOLDL, True)], 15, after=4)
runs(ftf, [("−  Σ (peso × fuerza) costo", GOLDL, True)], 15, after=4)
runs(ftf, [("−  Σ (peso × fuerza) riesgo ético", GOLDL, True)], 15, after=10)
par(ftf, "Determinista y auditable: mismos datos → mismo ranking.", 11.5, RGBColor(0xC7,0xD3,0xE3))

# ==========================================================================
# 6 — Fases 7-8
# ==========================================================================
s = slide(); header(s, "FASES 7–8", "Trazabilidad, bandera ética y tablero")
footer(s, 6)
items = [
    ("Trazabilidad de derivados", "Cuando se mitiga un requisito, nacen requisitos derivados que quedan "
     "ligados a su origen. Si un derivado obligatorio queda por debajo de su padre, salta una alerta."),
    ("Bandera ética", "Cada requisito muestra en el tablero una señal derivada de su análisis: verde, "
     "atención o crítico."),
    ("Tablero y exportación", "Ranking, desglose por categoría y bandera en una sola vista; todo el "
     "proyecto se puede exportar (JSON / CSV) para auditar."),
]
y = I(2.0)
for t,b in items:
    rect(s, I(0.7), y, I(11.9), I(1.35), CANVAS, line=GOLDL)
    rect(s, I(0.7), y, I(0.14), I(1.35), GOLD)
    ntf = textbox(s, I(1.05), y+I(0.16), I(11.3), I(1.05), MSO_ANCHOR.MIDDLE)
    par(ntf, t, 16, NAVY, True, first=True, after=3)
    par(ntf, b, 13.5, GREY)
    y += I(1.55)

# ==========================================================================
# 7 — Ejemplo: contexto
# ==========================================================================
s = slide(NAVY)
rect(s, 0, I(2.7), SW, Pt(3), GOLD)
tf = textbox(s, I(0.9), I(2.05), I(11.5), I(0.6))
par(tf, "UN EJEMPLO PARA VERLO CLARO", 14, GOLD, True, first=True)
tf2 = textbox(s, I(0.9), I(2.95), I(11.5), I(2.6))
par(tf2, "“MediTurno”", 34, WHITE, True, first=True, after=6)
par(tf2, "App que agenda y prioriza horas médicas con ayuda de IA.\n"
         "Analizaremos dos de sus requisitos con el método.", 18, GOLDL, False)
tf3 = textbox(s, I(0.9), I(5.8), I(11.5), I(1.0))
par(tf3, "Ejemplo ilustrativo — distinto al caso que ustedes trabajarán en la actividad.",
    13, RGBColor(0xB9,0xC6,0xD8), italic=True, first=True)

# ==========================================================================
# 8 — Ejemplo: análisis ético de los requisitos
# ==========================================================================
s = slide(); header(s, "EJEMPLO · FASES 2–3", "Qué detecta el análisis en cada requisito")
footer(s, 8)
reqs = [
    ("R1", "Recordatorio de citas por SMS", GREEN, "Riesgo bajo",
     "Beneficio claro, sin datos sensibles ni decisiones automáticas. Refuerza el acceso."),
    ("R2", "Priorizar horas con un “score de riesgo” que usa edad, comuna y previsión de salud", RED, "Riesgo alto",
     "Variables proxy → posible discriminación (equidad); usa datos de salud → privacidad. "
     "Norma tensionada: Ley 20.609 y Ley 19.628 / GDPR."),
]
y = I(1.75)
for code, name, badge, blab, desc in reqs:
    h = I(2.35)
    rect(s, I(0.7), y, I(11.9), h, CANVAS, line=GOLDL)
    rect(s, I(0.7), y, I(1.3), h, NAVY)
    ctf = textbox(s, I(0.7), y+I(0.2), I(1.3), h-I(0.4), MSO_ANCHOR.MIDDLE)
    par(ctf, code, 30, GOLD, True, align=PP_ALIGN.CENTER, first=True)
    # badge
    rect(s, I(10.9), y+I(0.25), I(1.5), I(0.4), badge)
    btf = textbox(s, I(10.9), y+I(0.25), I(1.5), I(0.4), MSO_ANCHOR.MIDDLE)
    par(btf, blab, 11, WHITE, True, align=PP_ALIGN.CENTER, first=True)
    ntf = textbox(s, I(2.2), y+I(0.28), I(8.5), h-I(0.5))
    par(ntf, name, 16, NAVY, True, first=True, after=6)
    par(ntf, desc, 13.5, GREY)
    y += h + I(0.25)

# ==========================================================================
# 9 — Ejemplo: priorización con números
# ==========================================================================
s = slide(); header(s, "EJEMPLO · FASES 4–6", "La priorización, con números")
footer(s, 9)
tf = textbox(s, I(0.7), I(1.55), I(11.9), I(0.55))
par(tf, "Cuatro dimensiones (con su peso). Cada celda es la fuerza 0–5 del requisito en esa dimensión.",
    14, INK, first=True)
# tabla
rows, cols = 3, 6
tbl_shape = s.shapes.add_table(rows, cols, I(0.7), I(2.25), I(11.9), I(1.7))
tbl = tbl_shape.table; style_table(tbl)
widths = [I(3.4), I(2.0), I(1.85), I(1.6), I(1.85), I(1.2)]
for j,w in enumerate(widths): tbl.columns[j].width = w
heads = ["Requisito", "Valor paciente\nbenef. (5)", "Equidad\nético (4)",
         "Costo\ncosto (2)", "Privacidad\nriesgo (4)", "Puntaje"]
for j,htext in enumerate(heads):
    set_cell(tbl.cell(0,j), htext, 11, WHITE, True, NAVY,
             PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER)
data = [
    ["R1 · Recordatorio SMS", "4", "3", "1", "1", "26"],
    ["R2 · Score de riesgo", "5", "1", "3", "5", "3"],
]
for i,row in enumerate(data, start=1):
    fill = WHITE if i==1 else CANVAS
    for j,val in enumerate(row):
        col = INK; bold = (j==0)
        align = PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER
        if j==5:
            bold = True; col = GREEN if i==1 else RED; fill = GOLDL
        set_cell(tbl.cell(i,j), val, 12, col, bold, fill, align)
# cálculo
tf = textbox(s, I(0.7), I(4.25), I(11.9), I(1.7))
runs(tf, [("R1 = (5×4)+(4×3)−(2×1)−(4×1) = 20+12−2−4 = ", INK, False), ("26", GREEN, True)], 15, first=True, after=6)
runs(tf, [("R2 = (5×5)+(4×1)−(2×3)−(4×5) = 25+4−6−20 = ", INK, False), ("3", RED, True)], 15, after=10)
runs(tf, [("Resultado: ", NAVY, True),
          ("R1 se prioriza sobre R2. El alto riesgo ético de R2 lo hace caer, pese a su utilidad.", INK, False)],
     15)
rect(s, I(0.7), I(6.35), I(11.9), Pt(3), GOLD)

# ==========================================================================
# 10 — Su rol en la experimentación
# ==========================================================================
s = slide(); header(s, "LO QUE HARÁN USTEDES", "Su rol en la sesión de prueba")
footer(s, 10)
pts = [
    ("Reciben los requisitos ya escritos", "No los definen: los analizan, tratan y priorizan con la app."),
    ("Recorren las fases 1 → 8", "Registro, análisis ético, priorización, trazabilidad y tablero."),
    ("Deciden el tratamiento", "En al menos uno o dos requisitos: aceptar, reformular, mitigar o eliminar."),
    ("Nos dan su mirada", "Qué fue útil, qué confundió y qué riesgos les sorprendió que detectara."),
]
y = I(1.9)
for i,(t,b) in enumerate(pts):
    rect(s, I(0.7), y, I(0.55), I(0.9), NAVY, shape=MSO_SHAPE.OVAL)
    ntf = textbox(s, I(0.7), y, I(0.55), I(0.9), MSO_ANCHOR.MIDDLE)
    par(ntf, str(i+1), 18, GOLD, True, align=PP_ALIGN.CENTER, first=True)
    ttf = textbox(s, I(1.5), y+I(0.02), I(11.0), I(0.95), MSO_ANCHOR.MIDDLE)
    par(ttf, t, 16, NAVY, True, first=True, after=2)
    par(ttf, b, 13.5, GREY)
    y += I(1.18)
tf = textbox(s, I(0.7), I(6.75), I(11.9), I(0.5))
par(tf, "No hay respuestas correctas ni se evalúa su desempeño: evaluamos la herramienta.", 13, INK, italic=True, first=True)

# ==========================================================================
# 11 — Cierre
# ==========================================================================
s = slide(NAVY)
rect(s, I(0.9), I(2.8), I(1.4), Pt(6), GOLD)
tf = textbox(s, I(0.9), I(3.0), I(11.5), I(2.0))
par(tf, "¿Comenzamos?", 40, WHITE, True, first=True, after=10)
par(tf, "Gracias por ayudarnos a probar y mejorar el método.", 18, GOLDL)
tf2 = textbox(s, I(0.9), I(6.4), I(11.5), I(0.5))
par(tf2, "Tesis · Nicolás Pérez   |   Universidad San Sebastián", 12, RGBColor(0xB9,0xC6,0xD8), first=True)

prs.save(OUT)
print("OK ->", OUT, "| slides:", len(prs.slides._sldIdLst))
